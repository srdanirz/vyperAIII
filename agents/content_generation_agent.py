# agents/content_generation_agent.py

import logging
from typing import Dict, Any, Optional, List
import json
from pathlib import Path
from datetime import datetime

# Document Generation
from docx import Document  # Cambio aquí
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

# Image handling
from PIL import Image
import matplotlib.pyplot as plt
import seaborn as sns

# Data processing
import pandas as pd
import numpy as np

from .base_agent import BaseAgent
from llm_factory import get_llm

logger = logging.getLogger(__name__)

class ContentStyle:
    """Definición de estilos y temas para contenido"""
    
    THEMES = {
        "professional": {
            "colors": {
                "primary": "#2C3E50",
                "secondary": "#E74C3C",
                "accent": "#3498DB",
                "background": "#FFFFFF",
                "text": "#2C3E50"
            },
            "fonts": {
                "title": "Calibri",
                "heading": "Calibri",
                "body": "Calibri"
            },
            "sizes": {
                "title": 44,
                "heading": 32,
                "subheading": 28,
                "body": 24
            }
        },
        "modern": {
            "colors": {
                "primary": "#1A237E",
                "secondary": "#FF4081",
                "accent": "#00BCD4",
                "background": "#FAFAFA",
                "text": "#212121"
            },
            "fonts": {
                "title": "Helvetica",
                "heading": "Helvetica",
                "body": "Helvetica"
            },
            "sizes": {
                "title": 40,
                "heading": 30,
                "subheading": 26,
                "body": 22
            }
        },
        "creative": {
            "colors": {
                "primary": "#6200EA",
                "secondary": "#00BFA5",
                "accent": "#FFD700",
                "background": "#FFFFFF",
                "text": "#424242"
            },
            "fonts": {
                "title": "Georgia",
                "heading": "Georgia",
                "body": "Arial"
            },
            "sizes": {
                "title": 42,
                "heading": 34,
                "subheading": 28,
                "body": 24
            }
        }
    }

    @staticmethod
    def get_theme(name: str = "professional") -> Dict[str, Any]:
        return ContentStyle.THEMES.get(name, ContentStyle.THEMES["professional"])
    
    @staticmethod
    def hex_to_rgb(hex_color: str) -> tuple:
        """Convierte color hex a RGB"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

class ContentGenerator:
    """Clase base para generadores de contenido específicos"""
    
    def __init__(self, theme: str = "professional"):
        self.theme = ContentStyle.get_theme(theme)
        self.output_dir = Path("output")
        self.output_dir.mkdir(exist_ok=True)
        
    async def generate(self, content: Dict[str, Any]) -> Dict[str, Any]:
        """Método base para generar contenido"""
        raise NotImplementedError

class PowerPointGenerator(ContentGenerator):
    """Generador específico para presentaciones PowerPoint"""
    
    def __init__(self, theme: str = "professional"):
        super().__init__(theme)
        self.prs = Presentation()
        
    def _apply_theme(self):
        """Aplica el tema a la presentación"""
        # Implementar aplicación de tema
        pass
        
    def _create_title_slide(self, title: str, subtitle: str = None):
        """Crea la diapositiva de título"""
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)
        
        # Configurar título
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(self.theme["sizes"]["title"])
        title_shape.text_frame.paragraphs[0].font.name = self.theme["fonts"]["title"]
        
        # Configurar subtítulo si existe
        if subtitle and slide.placeholders[1]:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(self.theme["sizes"]["subheading"])
            subtitle_shape.text_frame.paragraphs[0].font.name = self.theme["fonts"]["body"]

    def _create_content_slide(self, title: str, content: List[str], layout_index: int = 1):
        """Crea una diapositiva de contenido"""
        layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(layout)
        
        # Configurar título
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(self.theme["sizes"]["heading"])
        title_shape.text_frame.paragraphs[0].font.name = self.theme["fonts"]["heading"]
        
        # Configurar contenido
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        
        for idx, item in enumerate(content):
            p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(self.theme["sizes"]["body"])
            p.font.name = self.theme["fonts"]["body"]
            
    async def generate(self, content: Dict[str, Any]) -> Dict[str, Any]:
        """Genera la presentación PowerPoint"""
        try:
            self._apply_theme()
            
            # Crear diapositiva de título
            self._create_title_slide(
                content["title"],
                content.get("subtitle")
            )
            
            # Crear diapositivas de contenido
            for slide in content["slides"]:
                self._create_content_slide(
                    slide["title"],
                    slide["content"],
                    slide.get("layout_index", 1)
                )
            
            # Guardar presentación
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = self.output_dir / f"presentation_{timestamp}.pptx"
            self.prs.save(str(output_path))
            
            return {
                "type": "powerpoint",
                "path": str(output_path),
                "slides_count": len(self.prs.slides),
                "theme_used": self.theme
            }
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint: {e}")
            raise

class WordGenerator(ContentGenerator):
    """Generador específico para documentos Word"""
    
    def __init__(self, theme: str = "professional"):
        super().__init__(theme)
        self.doc = Document()
        
    def _apply_theme(self):
        """Aplica el tema al documento"""
        # Implementar estilos del documento
        pass
        
    def _add_title(self, title: str):
        """Añade el título principal del documento"""
        self.doc.add_heading(title, 0)
        
    def _add_heading(self, text: str, level: int):
        """Añade un encabezado al documento"""
        self.doc.add_heading(text, level)
        
    def _add_paragraph(self, text: str):
        """Añade un párrafo al documento"""
        self.doc.add_paragraph(text)
        
    async def generate(self, content: Dict[str, Any]) -> Dict[str, Any]:
        """Genera el documento Word"""
        try:
            self._apply_theme()
            
            # Añadir título principal
            self._add_title(content["title"])
            
            # Añadir contenido
            for section in content["sections"]:
                self._add_heading(section["title"], section.get("level", 1))
                for paragraph in section["content"]:
                    self._add_paragraph(paragraph)
            
            # Guardar documento
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = self.output_dir / f"document_{timestamp}.docx"
            self.doc.save(str(output_path))
            
            return {
                "type": "word",
                "path": str(output_path),
                "sections_count": len(content["sections"]),
                "theme_used": self.theme
            }
            
        except Exception as e:
            logger.error(f"Error generating Word document: {e}")
            raise

class PDFGenerator(ContentGenerator):
    """Generador específico para documentos PDF"""
    
    def __init__(self, theme: str = "professional"):
        super().__init__(theme)
        self.styles = getSampleStyleSheet()
        
    def _create_custom_styles(self):
        """Crea estilos personalizados para el PDF"""
        custom_styles = {}
        
        # Estilo de título
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=self.styles['Title'],
            fontSize=self.theme["sizes"]["title"],
            fontName=self.theme["fonts"]["title"],
            textColor=self.theme["colors"]["primary"]
        )
        custom_styles['Title'] = title_style
        
        # Estilo de encabezado
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=self.styles['Heading1'],
            fontSize=self.theme["sizes"]["heading"],
            fontName=self.theme["fonts"]["heading"],
            textColor=self.theme["colors"]["secondary"]
        )
        custom_styles['Heading'] = heading_style
        
        # Estilo de cuerpo
        body_style = ParagraphStyle(
            'CustomBody',
            parent=self.styles['Normal'],
            fontSize=self.theme["sizes"]["body"],
            fontName=self.theme["fonts"]["body"],
            textColor=self.theme["colors"]["text"]
        )
        custom_styles['Body'] = body_style
        
        return custom_styles
        
    async def generate(self, content: Dict[str, Any]) -> Dict[str, Any]:
        """Genera el documento PDF"""
        try:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = self.output_dir / f"document_{timestamp}.pdf"
            
            doc = SimpleDocTemplate(
                str(output_path),
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            # Crear estilos personalizados
            custom_styles = self._create_custom_styles()
            
            # Preparar elementos del documento
            elements = []
            
            # Añadir título
            elements.append(Paragraph(content["title"], custom_styles['Title']))
            elements.append(Spacer(1, 30))
            
            # Añadir contenido
            for section in content["sections"]:
                elements.append(Paragraph(section["title"], custom_styles['Heading']))
                elements.append(Spacer(1, 12))
                
                for paragraph in section["content"]:
                    elements.append(Paragraph(paragraph, custom_styles['Body']))
                    elements.append(Spacer(1, 12))
                
                elements.append(Spacer(1, 20))
            
            # Generar PDF
            doc.build(elements)
            
            return {
                "type": "pdf",
                "path": str(output_path),
                "sections_count": len(content["sections"]),
                "theme_used": self.theme
            }
            
        except Exception as e:
            logger.error(f"Error generating PDF: {e}")
            raise

class ContentGenerationAgent(BaseAgent):
    """
    Agente avanzado para generación de contenido con colaboración multiagente.
    """
    
    CONTENT_TYPES = {
        "powerpoint": PowerPointGenerator,
        "word": WordGenerator,
        "pdf": PDFGenerator
    }
    
    def __init__(
        self,
        task: str,
        openai_api_key: str,
        metadata: Optional[Dict[str, Any]] = None,
        shared_data: Optional[Dict[str, Any]] = None
    ):
        super().__init__(task, openai_api_key, metadata, shared_data)
        engine_mode = self.metadata.get("engine_mode", "openai")
        self.llm = get_llm(engine_mode, openai_api_key, model="gpt-4-turbo", temperature=0.7)
        
    def _determine_content_type(self) -> str:
        """Determina el tipo de contenido basado en la tarea"""
        task_lower = self.task.lower()
        
        # Mapeo de palabras clave a tipos de contenido
        content_mappings = {
            "powerpoint": ["powerpoint", "presentación", "slides", "diapositivas", "ppt"],
            "word": ["documento", "word", "texto", "doc", "docx"],
            "pdf": ["pdf", "documento pdf", "reporte pdf"]
        }
        
        for content_type, keywords in content_mappings.items():
            if any(keyword in task_lower for keyword in keywords):
                return content_type
                
        # Default a PowerPoint si no se especifica
        return "powerpoint"
        
    async def _generate_content_structure(self, content_type: str, research_data: str) -> Dict[str, Any]:
        """Genera la estructura del contenido usando el LLM"""
        try:
            messages = [
                {
                    "role": "system",
                    "content": f"Eres un experto en crear estructuras de {content_type}. "
                            "Debes generar una estructura profesional y detallada basada "
                            "en la investigación proporcionada."
                },
                {
                    "role": "user",
                    "content": (
                        f"Crea una estructura detallada para un {content_type} sobre: {self.task}\n\n"
                        f"Usando esta investigación:\n{research_data}\n\n"
                        "Responde en formato JSON con esta estructura:\n"
                        "{\n"
                        '  "title": "Título principal",\n'
                        '  "subtitle": "Subtítulo opcional",\n'
                        '  "theme": "professional",\n'
                        '  "slides": [\n'
                        '    {\n'
                        '      "title": "Título de slide",\n'
                        '      "content": ["Punto 1", "Punto 2"],\n'
                        '      "layout": "content",\n'
                        '      "notes": "Notas del presentador"\n'
                        '    }\n'
                        '  ]\n'
                        '}'
                    )
                }
            ]

            response = await self.llm.agenerate([messages])
            content = response.generations[0][0].message.content
            
            # Extraer el JSON del texto
            try:
                # Buscar el primer { y último }
                start = content.find('{')
                end = content.rfind('}') + 1
                if start >= 0 and end > start:
                    json_str = content[start:end]
                    return json.loads(json_str)
            except json.JSONDecodeError:
                logger.error("Error decodificando JSON de la respuesta")
                # Crear estructura por defecto
                return {
                    "title": "Presentación sobre IA",
                    "subtitle": "Conceptos y Aplicaciones",
                    "theme": "professional",
                    "slides": [
                        {
                            "title": "Introducción",
                            "content": ["Punto 1", "Punto 2"],
                            "layout": "content",
                            "notes": "Slide introductorio"
                        }
                    ]
                }

        except Exception as e:
            logger.error(f"Error generating content structure: {e}")
            raise

    async def _enhance_content_with_team(self, content: Dict[str, Any]) -> Dict[str, Any]:
        """Mejora el contenido usando feedback del equipo"""
        try:
            # Obtenemos análisis y sugerencias de otros agentes
            analysis_data = self.shared_data.get('analysis', {}).get('analysis_result', '')
            validation_data = self.shared_data.get('validation', {}).get('validation_summary', '')
            
            messages = [
                {
                    "role": "system",
                    "content": "Eres un experto en mejorar contenido basado en feedback. "
                              "Debes analizar el contenido actual y el feedback para sugerir mejoras."
                },
                {
                    "role": "user",
                    "content": f"""
                    Mejora este contenido:
                    {json.dumps(content, indent=2)}
                    
                    Basado en este análisis:
                    {analysis_data}
                    
                    Y esta validación:
                    {validation_data}
                    
                    Mantén el mismo formato JSON pero mejora:
                    - Títulos y subtítulos
                    - Puntos clave
                    - Estructura y organización
                    - Claridad y coherencia
                    """
                }
            ]
            
            resp = await self.llm.agenerate([messages])
            enhanced_content = json.loads(resp.generations[0][0].message.content)
            return enhanced_content
            
        except Exception as e:
            logger.error(f"Error enhancing content: {e}")
            raise

    async def _execute(self) -> Dict[str, Any]:
        """Ejecuta la generación de contenido"""
        try:
            # 1. Determinar tipo de contenido
            content_type = self._determine_content_type()
            logger.info(f"Generating {content_type} content")
            
            # 2. Obtener datos de investigación
            research_data = self.shared_data.get('research', {}).get('research_findings', '')
            
            # 3. Generar estructura inicial
            content_structure = await self._generate_content_structure(content_type, research_data)
            
            # 4. Mejorar contenido con feedback del equipo
            enhanced_content = await self._enhance_content_with_team(content_structure)
            
            # 5. Generar el contenido final
            generator_class = self.CONTENT_TYPES.get(content_type)
            if not generator_class:
                raise ValueError(f"No generator available for content type: {content_type}")
                
            generator = generator_class(theme=enhanced_content.get("theme", "professional"))
            result = await generator.generate(enhanced_content)
            
            # 6. Preparar respuesta
            return {
                "content_type": content_type,
                "output_file": result["path"],
                "metadata": {
                    "type": result["type"],
                    "theme": result["theme_used"],
                    "timestamp": datetime.now().isoformat(),
                    "content_structure": enhanced_content
                },
                "statistics": {
                    "slides_count": result.get("slides_count"),
                    "sections_count": result.get("sections_count")
                }
            }
            
        except Exception as e:
            logger.error(f"Error in ContentGenerationAgent: {e}")
            return {"error": str(e)}

    async def get_status(self) -> Dict[str, Any]:
        """Retorna el estado actual del agente"""
        return {
            "agent_type": "ContentGenerationAgent",
            "task": self.task,
            "metadata": self.metadata,
            "shared_data_keys": list(self.shared_data.keys())
        }